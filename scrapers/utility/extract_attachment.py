#!/usr/bin/env python3
"""
附件内容提取工具
支持格式：PDF, DOCX, XLSX, PPTX
用法：python3 extract_attachment.py <文件路径>
"""

import sys
import os

def extract_pdf(file_path):
    """提取 PDF 文件内容"""
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return '\n\n'.join(text_parts)
    except Exception as e:
        print(f"PDF 提取失败: {e}", file=sys.stderr)
        return ""

def extract_docx(file_path):
    """提取 DOCX 文件内容"""
    try:
        from docx import Document
        doc = Document(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # 提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip(' |'):
                    text_parts.append(row_text)
        
        return '\n\n'.join(text_parts)
    except Exception as e:
        print(f"DOCX 提取失败: {e}", file=sys.stderr)
        return ""

def extract_xlsx(file_path):
    """提取 XLSX 文件内容"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True)
        text_parts = []
        
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text_parts.append(f"=== 工作表: {sheet} ===")
            for row in ws.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip(' |'):
                    text_parts.append(row_text)
        
        wb.close()
        return '\n'.join(text_parts)
    except Exception as e:
        print(f"XLSX 提取失败: {e}", file=sys.stderr)
        return ""

def extract_pptx(file_path):
    """提取 PPTX 文件内容"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"=== 幻灯片 {slide_num} ==="]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            if len(slide_texts) > 1:
                text_parts.append('\n'.join(slide_texts))
        
        return '\n\n'.join(text_parts)
    except Exception as e:
        print(f"PPTX 提取失败: {e}", file=sys.stderr)
        return ""

def main():
    if len(sys.argv) != 2:
        print("用法: python3 extract_attachment.py <文件路径>", file=sys.stderr)
        sys.exit(1)
    
    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(f"文件不存在: {file_path}", file=sys.stderr)
        sys.exit(1)
    
    # 根据文件扩展名选择提取方法
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        content = extract_pdf(file_path)
    elif ext == '.docx' or ext == '.doc':
        content = extract_docx(file_path)
    elif ext == '.xlsx' or ext == '.xls':
        content = extract_xlsx(file_path)
    elif ext == '.pptx' or ext == '.ppt':
        content = extract_pptx(file_path)
    else:
        print(f"不支持的文件格式: {ext}", file=sys.stderr)
        sys.exit(1)
    
    if content:
        print(content)
    else:
        print("无法提取文件内容", file=sys.stderr)
        sys.exit(1)

if __name__ == '__main__':
    main()
